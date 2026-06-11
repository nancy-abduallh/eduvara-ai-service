"""
pipeline/video_assembler.py
============================
Assembles the final lecture video:
  1. Generate per-slide TTS audio via edge-tts
  2. Convert PPTX → PDF → PNG frames (LibreOffice + pdf2image)
  3. Render each slide as an MP4 with ffmpeg (zoom-pan + audio)
  4. Generate avatar lip-sync clips via Wav2Lip
     (checkpoint path comes from settings.WAV2LIP_CHECKPOINT_PATH,
      avatar image path comes from settings.AVATAR_STAND_PATH —
      same pattern as the script model being loaded from
      settings.SCRIPT_ADAPTER_PATH in model_registry.py)
  5. Merge intro + lecture + outro into the final video
"""

from __future__ import annotations

import asyncio
import logging
import os
import random
import subprocess
from pathlib import Path
from typing import List, Optional

logger = logging.getLogger("edugenie.video")

TRANSITIONS = ["slideleft", "slideright", "wipeleft", "circleopen", "fade"]
VOICE       = "en-US-AvaNeural"
FADE_DUR    = 0.5


# ── ffmpeg / ffprobe path resolution ─────────────────────────────────────────

def _find_ffmpeg_tool(name: str) -> str:
    """
    Return the full path to an ffmpeg-suite binary (ffmpeg or ffprobe).

    Search order:
      1. Already on PATH  (Linux, macOS, Windows with ffmpeg in PATH)
      2. Common Windows install locations:
         - winget/manual:  C:\\ffmpeg\\bin\\<name>.exe
         - scoop:          C:\\Users\\<user>\\scoop\\apps\\ffmpeg\\current\\bin\\<name>.exe
         - chocolatey:     C:\\ProgramData\\chocolatey\\bin\\<name>.exe
         - alongside this file (portable drop-in)

    Raises RuntimeError with install instructions if not found anywhere.
    """
    import shutil as _shutil
    import sys

    # 1. PATH (works on any OS when ffmpeg is properly installed)
    found = _shutil.which(name)
    if found:
        return found

    # 2. Common Windows locations
    if sys.platform == "win32":
        import glob
        import getpass
        username = getpass.getuser()
        candidates = [
            rf"C:\ffmpeg\bin\{name}.exe",
            rf"C:\Program Files\ffmpeg\bin\{name}.exe",
            rf"C:\Program Files (x86)\ffmpeg\bin\{name}.exe",
            rf"C:\Users\{username}\scoop\apps\ffmpeg\current\bin\{name}.exe",
            rf"C:\ProgramData\chocolatey\bin\{name}.exe",
            # Portable: drop ffmpeg/bin/ next to the pipeline folder
            str(Path(__file__).resolve().parents[1] / "ffmpeg" / "bin" / f"{name}.exe"),
            str(Path(__file__).resolve().parent / "ffmpeg" / "bin" / f"{name}.exe"),
        ]
        # Also handle versioned scoop paths like ffmpeg\7.1\bin\
        candidates += glob.glob(
            rf"C:\Users\{username}\scoop\apps\ffmpeg\*\bin\{name}.exe"
        )
        for path in candidates:
            if os.path.isfile(path):
                logger.info("Found %s at: %s", name, path)
                return path

    raise RuntimeError(
        f"\n{'='*60}\n"
        f"  {name} not found on this system.\n\n"
        f"  ── How to install ffmpeg on Windows (pick ONE) ──\n\n"
        f"  Option A — Winget (built into Windows 10/11):\n"
        f"    1. Open PowerShell as Administrator\n"
        f"    2. winget install --id Gyan.FFmpeg -e --source winget\n"
        f"    3. Restart your terminal / uvicorn\n\n"
        f"  Option B — Manual (portable, no admin needed):\n"
        f"    1. Download from https://www.gyan.dev/ffmpeg/builds/\n"
        f"       → ffmpeg-release-essentials.zip\n"
        f"    2. Extract to C:\\ffmpeg\\\n"
        f"       (so you have C:\\ffmpeg\\bin\\ffmpeg.exe)\n"
        f"    3. Add C:\\ffmpeg\\bin to your System PATH:\n"
        f"       Search 'Edit the system environment variables' → PATH → New\n"
        f"    4. Restart your terminal / uvicorn\n\n"
        f"  Option C — Scoop:\n"
        f"    scoop install ffmpeg\n\n"
        f"  Option D — Chocolatey:\n"
        f"    choco install ffmpeg\n"
        f"{'='*60}"
    )


# Resolve once at import time so every function uses the same cached path.
# We wrap in a lambda so the error is raised lazily (at first use), not at import.
def _ffmpeg()  -> str: return _find_ffmpeg_tool("ffmpeg")
def _ffprobe() -> str: return _find_ffmpeg_tool("ffprobe")


# ── TTS ───────────────────────────────────────────────────────────────────────

async def _tts(text: str, voice: str, out_path: str):
    import edge_tts
    await edge_tts.Communicate(text, voice).save(out_path)


async def generate_slide_audios(speaker_notes: List[str], workdir: str) -> List[str]:
    """
    Generate one MP3 per speaker note.
    Returns list of audio file paths (same length as speaker_notes).
    """
    paths = []
    for i, note in enumerate(speaker_notes):
        if not note or not note.strip():
            paths.append(None)
            continue
        out = os.path.join(workdir, f"slide_{i + 1}.mp3")
        try:
            await _tts(note, VOICE, out)
            paths.append(out)
        except Exception as e:
            logger.error(f"TTS failed for slide {i + 1}: {e}")
            paths.append(None)
    return paths


# ── PPTX → PNG frames ─────────────────────────────────────────────────────────

def _find_libreoffice() -> Optional[str]:
    """
    Return the LibreOffice executable path, or None if not found.
    Checks PATH first, then common Windows install locations including
    versioned directories (e.g. 'LibreOffice 26.2', 'LibreOffice 7.6').
    """
    import shutil as _shutil
    import sys

    # 1. Already on PATH (Linux / macOS / Windows with soffice in PATH)
    for cmd in ("libreoffice", "soffice"):
        found = _shutil.which(cmd)
        if found:
            return found

    # 2. Common Windows installation directories
    if sys.platform == "win32":
        import glob
        patterns = [
            # Standard install (version-agnostic directory name)
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            # Versioned directory names — LibreOffice 7.x, 24.x, 26.x …
            r"C:\Program Files\LibreOffice*\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice*\program\soffice.exe",
            # Portable / side-by-side installs
            r"C:\LibreOffice*\program\soffice.exe",
        ]
        for pattern in patterns:
            matches = sorted(glob.glob(pattern), reverse=True)  # newest first
            if matches:
                logger.info("Found LibreOffice at: %s", matches[0])
                return matches[0]

    return None


def _pptx_to_frames_python(pptx_path: str, workdir: str) -> List[str]:
    """
    Pure-Python fallback: render each PPTX slide to a 1920×1080 PNG using
    python-pptx + Pillow.  No LibreOffice or poppler required.

    Rendering quality is lower than LibreOffice (no fonts, no EMF shapes),
    but the job always completes and produces a watchable video.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from PIL import Image, ImageDraw, ImageFont
        import textwrap
    except ImportError as e:
        raise RuntimeError(
            f"python-pptx or Pillow not installed ({e}). "
            "Run: pip install python-pptx Pillow"
        )

    W, H = 1920, 1080
    prs   = Presentation(pptx_path)
    paths = []

    # Try to load a readable font; fall back to Pillow's default.
    FONT_PATHS = [
        "C:/Windows/Fonts/calibri.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ]
    title_font = body_font = None
    for fp in FONT_PATHS:
        if os.path.exists(fp):
            try:
                title_font = ImageFont.truetype(fp, 54)
                body_font  = ImageFont.truetype(fp, 30)
                break
            except Exception:
                pass
    if title_font is None:
        title_font = ImageFont.load_default()
        body_font  = ImageFont.load_default()

    for idx, slide in enumerate(prs.slides):
        img  = Image.new("RGB", (W, H), color=(18, 18, 28))   # dark background
        draw = ImageDraw.Draw(img)

        # Accent bar at top
        draw.rectangle([0, 0, W, 8], fill=(0, 188, 212))

        title_text = ""
        body_lines: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            raw = shape.text_frame.text.strip()
            if not raw:
                continue
            # Heuristic: placeholder 0 / 1 = title; rest = body
            ph_idx = getattr(getattr(shape, "placeholder_format", None), "idx", None)
            if ph_idx in (0, 1) and not title_text:
                title_text = raw
            else:
                body_lines.extend(raw.splitlines())

        # Draw title
        if title_text:
            draw.text((80, 60), title_text, font=title_font, fill=(0, 188, 212))
            draw.line([(80, 130), (W - 80, 130)], fill=(0, 188, 212, 80), width=2)

        # Draw body — word-wrap each line to ~95 chars
        y = 160
        for line in body_lines:
            for wrapped in textwrap.wrap(line, width=95) or [""]:
                if y > H - 60:
                    break
                draw.text((80, y), wrapped, font=body_font, fill=(230, 230, 230))
                y += 38

        # Slide number
        draw.text(
            (W - 80, H - 50),
            f"{idx + 1}/{len(prs.slides)}",
            font=body_font,
            fill=(120, 120, 120),
        )

        p = os.path.join(workdir, f"slide_{idx + 1}.png")
        img.save(p, "PNG")
        paths.append(p)
        logger.debug("Fallback renderer: saved %s", p)

    logger.warning(
        "Used pure-Python PPTX renderer (no LibreOffice). "
        "Install LibreOffice and add it to PATH for higher-fidelity slides."
    )
    return paths


def _find_poppler_path() -> Optional[str]:
    """
    Return the poppler bin/ directory path on Windows, or None on other platforms.

    pdf2image needs poppler's pdftoppm.exe in PATH (or passed as poppler_path).
    If poppler is not on PATH, check common Windows install locations.

    Install poppler on Windows (one-time, no admin needed):
      1. Download from https://github.com/oschwartz10612/poppler-windows/releases
         (latest release ZIP, e.g. Release-24.08.0-0.zip)
      2. Extract to C:\\poppler\\ so you have C:\\poppler\\Library\\bin\\pdftoppm.exe
      3. OR add the bin/ directory to your System PATH and restart the terminal.
    """
    import shutil as _shutil
    import sys

    # 1. Already on PATH
    if _shutil.which("pdftoppm"):
        return None  # None = let pdf2image find it via PATH (default behaviour)

    if sys.platform != "win32":
        return None  # On Linux/macOS, assume PATH is correct

    import glob
    import getpass
    username = getpass.getuser()
    candidates = [
        # Manual / portable extraction
        r"C:\poppler\Library\bin",
        r"C:\poppler\bin",
        r"C:\Program Files\poppler\bin",
        r"C:\Program Files (x86)\poppler\bin",
        # Alongside the project
        str(Path(__file__).resolve().parents[1] / "poppler" / "bin"),
        str(Path(__file__).resolve().parents[1] / "poppler" / "Library" / "bin"),
        # Scoop
        rf"C:\Users\{username}\scoop\apps\poppler\current\bin",
        # Chocolatey
        rf"C:\ProgramData\chocolatey\bin",
    ]
    # Also handle versioned scoop paths
    candidates += glob.glob(rf"C:\Users\{username}\scoop\apps\poppler\*\bin")
    # Glob for any C:\poppler-* directory (versioned releases)
    candidates += glob.glob(r"C:\poppler-*\Library\bin")
    candidates += glob.glob(r"C:\poppler-*\bin")

    for path in candidates:
        if os.path.isfile(os.path.join(path, "pdftoppm.exe")):
            logger.info("Found poppler at: %s", path)
            return path

    logger.warning(
        "poppler not found — pdf2image will fail. "
        "To fix: download poppler for Windows from "
        "https://github.com/oschwartz10612/poppler-windows/releases, "
        "extract to C:\\poppler\\, then either add C:\\poppler\\Library\\bin to "
        "your System PATH or place the folder next to the project as poppler/."
    )
    return None


def _soffice_env() -> dict:
    """
    Return a clean subprocess environment for LibreOffice.

    LibreOffice ships its own Python interpreter under program/python-core-*/
    and adds that directory to PATH and PYTHONPATH when soffice starts.
    If our uvicorn process was started while soffice was already on the system
    PATH, LibreOffice's program/ directory may end up on sys.path, which causes
    'Cannot find module pydantic_settings' because LibreOffice's Python knows
    nothing about our venv.

    The fix: strip PYTHONPATH, PYTHONHOME, and LibreOffice's program/ directory
    from PATH before spawning soffice so its Python stays completely isolated.
    """
    import sys
    env = os.environ.copy()

    # Remove env vars that bleed our venv's Python settings into LO's process
    for var in ("PYTHONPATH", "PYTHONHOME", "PYTHONSTARTUP", "PYTHONUSERBASE"):
        env.pop(var, None)

    # Strip LibreOffice's own program/ directory from PATH so it cannot be
    # imported by any child process that tries to find Python modules.
    if sys.platform == "win32" and "PATH" in env:
        parts = env["PATH"].split(os.pathsep)
        cleaned = [
            p for p in parts
            if "LibreOffice" not in p or "program" not in p.lower()
        ]
        env["PATH"] = os.pathsep.join(cleaned)

    return env


def _pdf_to_pngs_pypdfium2(pdf_path: str, workdir: str, dpi: int = 144) -> Optional[List[str]]:
    """
    Render every page of a PDF to a PNG using pypdfium2.

    pypdfium2 bundles its own pdfium binary — no poppler, no system dependency.
    Install once: pip install pypdfium2

    Returns a list of PNG paths (one per page) or None if pypdfium2 is not installed.
    """
    try:
        import pypdfium2 as pdfium
    except ImportError:
        logger.debug("pypdfium2 not installed — skipping. Run: pip install pypdfium2")
        return None

    try:
        pdf   = pdfium.PdfDocument(pdf_path)
        scale = dpi / 72.0          # pdfium renders at 72 DPI base
        paths: List[str] = []
        for i in range(len(pdf)):
            page    = pdf[i]
            bitmap  = page.render(scale=scale, rotation=0)
            pil_img = bitmap.to_pil()
            p = os.path.join(workdir, f"slide_{i + 1}.png")
            pil_img.save(p, "PNG")
            paths.append(p)
        logger.info("pptx_to_frames: %d frames via pypdfium2 (no poppler)", len(paths))
        return paths
    except Exception as exc:
        logger.warning("pypdfium2 rendering failed (%s) — trying next method.", exc)
        return None


def _pptx_to_frames_libreoffice(pptx_path: str, workdir: str, soffice: str) -> Optional[List[str]]:
    """
    Export every PPTX slide to PNG via LibreOffice.

    Strategy inside this function (in order):
      A. LO --convert-to "png:impress_png_Export"
         Uses the named Impress export filter which exports ALL slides on
         LibreOffice 7.x and newer.  Each slide becomes <stem><N>.png.
      B. LO --convert-to pdf  →  pypdfium2 PNG render
         Guaranteed to export all pages.  pypdfium2 bundles pdfium — no poppler.
      C. LO --convert-to pdf  →  pdf2image (needs poppler)
         Legacy path kept for completeness.

    All subprocess calls use _soffice_env() to strip LibreOffice's Python from
    the environment so it cannot shadow our venv packages.
    """
    import glob as _glob
    stem    = Path(pptx_path).stem
    lo_env  = _soffice_env()

    # ── A: LO Impress PNG filter (exports ALL slides) ─────────────────────────
    # The filter string "png:impress_png_Export" tells LO to use the named
    # Impress export filter instead of the generic image converter, which on
    # many LO versions only converts the first slide.
    for filter_str in ("png:impress_png_Export", "png"):
        try:
            result = subprocess.run(
                [
                    soffice, "--headless", "--norestore",
                    "--convert-to", filter_str,
                    "--outdir", workdir,
                    pptx_path,
                ],
                capture_output=True, text=True, timeout=180,
                env=lo_env,
            )
        except subprocess.TimeoutExpired:
            logger.warning("LibreOffice PNG export timed out (filter=%s).", filter_str)
            continue
        except Exception as exc:
            logger.warning("LibreOffice PNG export error (filter=%s): %s", filter_str, exc)
            continue

        # Collect all PNGs that match <stem>*.png in the workdir (sorted)
        found = sorted(_glob.glob(os.path.join(workdir, f"{stem}*.png")))

        # Only accept if we got MORE than 1 file, or exactly 1 when the PPTX
        # genuinely has 1 slide.  If LO produced exactly 1 PNG for a multi-slide
        # PPTX it only converted the first slide — fall through.
        if found:
            try:
                from pptx import Presentation as _Prs
                expected = len(_Prs(pptx_path).slides)
            except Exception:
                expected = 0

            if len(found) >= max(1, expected):
                # Rename to canonical slide_N.png
                renamed: List[str] = []
                for i, src in enumerate(found):
                    dst = os.path.join(workdir, f"slide_{i + 1}.png")
                    # Avoid renaming if it would overwrite an existing slide_N.png
                    if src != dst:
                        try:
                            os.replace(src, dst)
                        except Exception:
                            dst = src
                    renamed.append(dst)
                logger.info(
                    "pptx_to_frames: %d/%d slides via LibreOffice PNG filter '%s'",
                    len(renamed), expected, filter_str,
                )
                return renamed
            else:
                logger.warning(
                    "LO PNG filter '%s' produced %d file(s) but PPTX has %d slides "
                    "— falling through to PDF path.",
                    filter_str, len(found), expected,
                )
                # Remove the partial output so it doesn't confuse the PDF path
                for f in found:
                    try:
                        os.remove(f)
                    except OSError:
                        pass

    # ── B/C: LO → PDF → PNG render ────────────────────────────────────────────
    # PDF export always produces one file containing ALL pages.
    logger.info("Falling back to LibreOffice PDF export path…")
    try:
        pdf_result = subprocess.run(
            [soffice, "--headless", "--norestore",
             "--convert-to", "pdf", "--outdir", workdir, pptx_path],
            capture_output=True, text=True, timeout=120,
            env=lo_env,
        )
    except Exception as exc:
        logger.warning("LibreOffice PDF export failed: %s", exc)
        return None

    pdf_path = os.path.join(workdir, stem + ".pdf")
    if not os.path.exists(pdf_path):
        logger.warning(
            "LibreOffice PDF export produced no file (rc=%d). stderr: %s",
            pdf_result.returncode, pdf_result.stderr[:300],
        )
        return None

    # B: pypdfium2 (no poppler needed)
    frames = _pdf_to_pngs_pypdfium2(pdf_path, workdir)
    if frames:
        return frames

    # C: pdf2image (needs poppler — legacy)
    try:
        from pdf2image import convert_from_path
        poppler_path = _find_poppler_path()
        kwargs: dict = {"dpi": 144}
        if poppler_path:
            kwargs["poppler_path"] = poppler_path
        images = convert_from_path(pdf_path, **kwargs)
        paths: List[str] = []
        for i, img in enumerate(images):
            p = os.path.join(workdir, f"slide_{i + 1}.png")
            img.save(p, "PNG")
            paths.append(p)
        logger.info("pptx_to_frames: %d frames via LibreOffice+pdf2image", len(paths))
        return paths
    except ImportError:
        logger.warning("pdf2image not installed (run: pip install pdf2image) — skipping.")
    except Exception as exc:
        logger.warning("pdf2image conversion failed (%s) — falling back.", exc)

    return None


def pptx_to_frames(pptx_path: str, workdir: str) -> List[str]:
    """
    Convert PPTX to a list of PNG paths (one per slide), ready for ffmpeg.

    Strategy (in order):
    1. LibreOffice: PNG filter  →  PDF+pypdfium2  →  PDF+pdf2image
       (best quality; LibreOffice renders fonts and shapes correctly)
    2. Pure-Python python-pptx + Pillow
       (always works, lower fidelity — no LibreOffice or PDF tools required)
    """
    soffice = _find_libreoffice()

    if soffice:
        frames = _pptx_to_frames_libreoffice(pptx_path, workdir, soffice)
        if frames:
            return frames
        logger.warning(
            "All LibreOffice paths failed — falling back to pure-Python renderer."
        )
    else:
        logger.warning(
            "LibreOffice not found — using pure-Python renderer fallback. "
            "Install LibreOffice from https://www.libreoffice.org/download/download/ "
            "for higher quality slides."
        )

    return _pptx_to_frames_python(pptx_path, workdir)


# ── ffprobe helpers ───────────────────────────────────────────────────────────

def _get_duration(path: str) -> float:
    r = subprocess.run(
        [_ffprobe(), "-v", "quiet", "-show_entries", "format=duration",
         "-of", "default=noprint_wrappers=1:nokey=1", path],
        capture_output=True, text=True,
    )
    return float(r.stdout.strip())


def _probe_dim(path: str) -> tuple[int, int]:
    r = subprocess.run(
        [_ffprobe(), "-v", "error", "-select_streams", "v:0",
         "-show_entries", "stream=width,height",
         "-of", "csv=p=0", path],
        capture_output=True, text=True,
    )
    parts = r.stdout.strip().split(",")
    return int(parts[0]), int(parts[1])


# ── Per-slide video render ────────────────────────────────────────────────────

def _find_windows_font() -> Optional[str]:
    """
    Return an absolute path to a TTF/OTF font file on Windows (or None elsewhere).

    Passing fontfile= to ffmpeg's drawtext filter bypasses fontconfig entirely —
    no fonts.conf lookup, no 'Fontconfig error: Cannot load default config file'.
    This is the correct fix for portable ffmpeg builds on Windows.
    """
    import sys
    if sys.platform != "win32":
        return None
    candidates = [
        r"C:\Windows\Fonts\calibrib.ttf",   # Calibri Bold
        r"C:\Windows\Fonts\calibri.ttf",    # Calibri Regular
        r"C:\Windows\Fonts\arialbd.ttf",    # Arial Bold
        r"C:\Windows\Fonts\arial.ttf",      # Arial Regular
        r"C:\Windows\Fonts\segoeui.ttf",    # Segoe UI
        r"C:\Windows\Fonts\tahoma.ttf",     # Tahoma
        r"C:\Windows\Fonts\verdana.ttf",    # Verdana
    ]
    for p in candidates:
        if os.path.isfile(p):
            return p
    return None


# Resolved once at module level — used by _build_vf_filter.
_WINDOWS_FONT: Optional[str] = _find_windows_font()


def _build_vf_filter(bar_width: int, slide_label: str, with_drawtext: bool = True) -> str:
    """
    Build the -vf filter string for a slide.

    drawtext with fontconfig fails on Windows with portable ffmpeg builds
    (no fonts.conf → "Fontconfig error: Cannot load default config file").

    FIX: When a Windows system font is found, pass fontfile= directly to
    drawtext — this completely bypasses fontconfig and always works.
    Fall back to with_drawtext=False (drawbox only) only when no font is found.
    """
    entrance_filter = (
        "zoompan="
        "z='if(lte(on,30),1.03-0.001*on,1.0)':"
        "x='iw/2-(iw/zoom/2)':"
        "y='ih/2-(ih/zoom/2)':"
        "d=1:s=1920x1080:fps=24,"
        "fade=t=in:st=0:d=0.5"
    )
    base = (
        "scale=1920:1080:force_original_aspect_ratio=decrease,"
        "pad=1920:1080:(ow-iw)/2:(oh-ih)/2:color=black,"
        "setsar=1,"
        f"{entrance_filter},"
        f"drawbox=x=0:y=1072:w={bar_width}:h=8:color=00BCD4@0.9:t=fill"
    )
    if with_drawtext:
        # Escape the label for ffmpeg filter syntax (colon and backslash are special)
        label_escaped = (
            slide_label
            .replace("\\", "\\\\")
            .replace(":", "\\:")
            .replace("'", "\\'")
        )
        if _WINDOWS_FONT:
            # Use fontfile= — bypasses fontconfig completely on Windows
            font_path_escaped = _WINDOWS_FONT.replace("\\", "/").replace(":", "\\:")
            drawtext = (
                f",drawtext=text='{label_escaped}':"
                f"fontfile='{font_path_escaped}':"
                f"fontsize=22:fontcolor=white@0.85:x=w-84:y=16:"
                "shadowx=1:shadowy=1:shadowcolor=black@0.6"
            )
        else:
            # Linux/macOS — fontconfig is available, use original syntax
            drawtext = (
                f",drawtext=text='{label_escaped}':fontsize=22:"
                "fontcolor=white@0.85:x=w-84:y=16:"
                "shadowx=1:shadowy=1:shadowcolor=black@0.6"
            )
        return base + drawtext
    return base


def _build_render_cmd(
    ffmpeg_bin: str,
    image_path: str,
    audio_path: Optional[str],
    output_path: str,
    vf_filter: str,
) -> list:
    """Return the ffmpeg command list for one slide render."""
    if audio_path and os.path.exists(audio_path):
        return [
            ffmpeg_bin, "-hide_banner", "-y",
            "-loop", "1", "-framerate", "24", "-i", image_path,
            "-i", audio_path,
            "-c:v", "libx264", "-b:v", "3000k", "-bufsize", "6000k",
            "-crf", "18", "-preset", "slow", "-tune", "stillimage",
            "-c:a", "aac", "-b:a", "192k", "-ar", "44100",
            "-pix_fmt", "yuv420p",
            "-vf", vf_filter,
            "-r", "24", "-shortest",
            output_path,
        ]
    # Silent 5-second slide — MUST include a silent AAC audio stream so that
    # concat_slide_videos can always reference [i:a] in its filtergraph.
    # We synthesise silence with the lavfi 'anullsrc' source, trimmed to 5 s.
    return [
        ffmpeg_bin, "-hide_banner", "-y",
        "-loop", "1", "-framerate", "24", "-t", "5", "-i", image_path,
        "-f", "lavfi", "-t", "5", "-i", "anullsrc=channel_layout=mono:sample_rate=44100",
        "-c:v", "libx264", "-b:v", "3000k", "-crf", "18",
        "-pix_fmt", "yuv420p",
        "-vf", vf_filter,
        "-r", "24",
        "-c:a", "aac", "-b:a", "192k", "-ar", "44100",
        "-shortest",
        output_path,
    ]


def render_slide_video(
    image_path:   str,
    audio_path:   Optional[str],
    output_path:  str,
    slide_index:  int,
    total_slides: int,
) -> bool:
    """
    Render one slide image + audio into an MP4 with entrance animation.

    Strategy:
      1. Try with full vf_filter (drawtext + drawbox + zoom/fade).
         On Windows, _build_vf_filter now uses fontfile= pointing at a Windows
         system font (e.g. Calibri/Arial), which completely bypasses fontconfig.
         This means attempt 0 should succeed on Windows without any retry needed.
      2. If ffmpeg still fails with a Fontconfig error (e.g. custom/portable
         builds without any system font), retry WITHOUT drawtext — only the
         progress bar + entrance animation remain.  This always succeeds because
         drawbox has no font dependency.
      3. Log an error and return False only if both attempts fail.
    """
    bar_width   = max(1, int(1920 * slide_index / max(1, total_slides)))
    slide_label = f"{slide_index}/{total_slides}"
    ff          = _ffmpeg()

    for attempt, with_drawtext in enumerate([True, False]):
        # Clean up any partial output from a previous attempt
        if attempt > 0 and os.path.exists(output_path):
            try:
                os.remove(output_path)
            except OSError:
                pass

        vf  = _build_vf_filter(bar_width, slide_label, with_drawtext=with_drawtext)
        cmd = _build_render_cmd(ff, image_path, audio_path, output_path, vf)

        result = subprocess.run(cmd, capture_output=True, text=True)
        success = (
            result.returncode == 0
            and os.path.exists(output_path)
            and os.path.getsize(output_path) > 0
        )

        if success:
            if attempt > 0:
                logger.info(
                    "Slide %d rendered OK (drawtext disabled — fontconfig unavailable).",
                    slide_index,
                )
            return True

        stderr_tail = result.stderr[-600:]

        # Detect fontconfig / drawtext failures and retry without drawtext
        is_fontconfig_error = (
            "Fontconfig error" in result.stderr
            or "Cannot load default config file" in result.stderr
            or ("drawtext" in result.stderr and "No such file" in result.stderr)
        )
        if attempt == 0 and is_fontconfig_error:
            logger.warning(
                "Slide %d: fontconfig error detected — retrying without drawtext overlay.",
                slide_index,
            )
            continue  # retry with with_drawtext=False

        # Not a fontconfig issue, or we already retried — log and give up
        logger.error("ffmpeg error on slide %d: %s", slide_index, stderr_tail)
        return False

    # Should be unreachable, but be safe
    logger.error("ffmpeg failed for slide %d after all attempts.", slide_index)
    return False


# ── Slide concatenation ───────────────────────────────────────────────────────

def concat_slide_videos(segment_paths: List[str], durations: List[float], output_path: str):
    """Merge all slide segments with random xfade transitions."""
    n = len(segment_paths)
    if n == 0:
        raise ValueError("No slide segments to concatenate")
    if n == 1:
        subprocess.run(
            [_ffmpeg(), "-y", "-i", segment_paths[0], "-c", "copy", output_path],
            check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        return

    transitions = [random.choice(TRANSITIONS) for _ in range(n)]
    inputs = []
    for v in segment_paths:
        inputs.extend(["-i", v])

    offset       = durations[0] - FADE_DUR
    filter_parts = []
    prev_v, prev_a = "[0:v]", "[0:a]"

    for i in range(1, n):
        trans = transitions[i]
        out_v = f"[v{i}]" if i < n - 1 else "[vout]"
        out_a = f"[a{i}]" if i < n - 1 else "[aout]"
        filter_parts.append(
            f"{prev_v}[{i}:v]xfade=transition={trans}:duration={FADE_DUR}:offset={offset:.3f}{out_v}"
        )
        filter_parts.append(
            f"{prev_a}[{i}:a]acrossfade=d={FADE_DUR}{out_a}"
        )
        prev_v  = out_v
        prev_a  = out_a
        offset += durations[i] - FADE_DUR

    cmd = (
        [_ffmpeg(), "-y"] + inputs +
        ["-filter_complex", "; ".join(filter_parts),
         "-map", "[vout]", "-map", "[aout]",
         "-c:v", "libx264", "-b:v", "3000k", "-crf", "18", "-preset", "slow",
         "-c:a", "aac", "-b:a", "192k", "-ar", "44100",
         "-pix_fmt", "yuv420p", "-r", "24",
         output_path]
    )
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(f"ffmpeg concat failed: {result.stderr[-800:]}")


# ── Wav2Lip avatar generation ─────────────────────────────────────────────────

def generate_avatar_clip(
    face_image:         str,
    audio_file:         str,
    output_file:        str,
    wav2lip_checkpoint: str,          # required — pass settings.WAV2LIP_CHECKPOINT_PATH
) -> Optional[str]:
    """
    Run Wav2Lip inference to produce a talking-head video.

    The checkpoint path is passed explicitly from settings (no hardcoded default)
    so the video pipeline uses the same config-driven pattern as the script
    generation pipeline, which reads its model path from settings.SCRIPT_ADAPTER_PATH.

    Returns the output_file path on success, None on failure.
    """
    if not wav2lip_checkpoint or not os.path.exists(wav2lip_checkpoint):
        logger.warning(
            f"Wav2Lip checkpoint not found at '{wav2lip_checkpoint}' — "
            "skipping avatar generation. Set WAV2LIP_CHECKPOINT_PATH in .env or config."
        )
        return None

    if not os.path.exists(face_image):
        logger.warning(f"Avatar face image not found at '{face_image}' — skipping.")
        return None

    if not os.path.exists(audio_file):
        logger.warning(f"Audio file not found at '{audio_file}' — skipping avatar clip.")
        return None

    result = subprocess.run(
        [
            "python", "inference.py",
            "--checkpoint_path", wav2lip_checkpoint,
            "--face",            face_image,
            "--audio",           audio_file,
            "--outfile",         output_file,
        ],
        capture_output=True, text=True,
    )
    if result.returncode != 0:
        logger.error(f"Wav2Lip failed: {result.stderr[-500:]}")
        return None

    if not os.path.exists(output_file) or os.path.getsize(output_file) == 0:
        logger.error(f"Wav2Lip produced no output file at {output_file}")
        return None

    logger.info(f"✅ Avatar clip saved: {output_file}")
    return output_file


# ── PiP (avatar overlay onto slides) ─────────────────────────────────────────

def overlay_avatar_on_slides(
    slides_video:  str,
    avatar_video:  str,
    output_path:   str,
    watermark:     str = "EduGenie | Adaptive Learning",
):
    """Overlay the full-body avatar in the bottom-right corner (PiP)."""
    av_src_w, av_src_h = _probe_dim(avatar_video)
    AV_DISP_H = 280
    AV_DISP_W = max(120, int(av_src_w / av_src_h * AV_DISP_H // 2 * 2))
    PAD = 27

    wm_escaped = watermark.replace("\\", "\\\\").replace(":", "\\:").replace("'", "\\'")
    # Use fontfile= on Windows to bypass fontconfig (same fix as _build_vf_filter)
    if _WINDOWS_FONT:
        font_path_escaped = _WINDOWS_FONT.replace("\\", "/").replace(":", "\\:")
        wm_filter = (
            f"drawtext=text='{wm_escaped}':"
            f"fontfile='{font_path_escaped}':"
            f"fontsize=20:fontcolor=white@0.55:"
            f"x=24:y=H-40:"
            f"shadowx=1:shadowy=1:shadowcolor=black@0.45"
        )
    else:
        wm_filter = (
            f"drawtext=text='{wm_escaped}':"
            f"fontsize=20:fontcolor=white@0.55:"
            f"x=24:y=H-40:"
            f"shadowx=1:shadowy=1:shadowcolor=black@0.45"
        )
    pip_filter = (
        f"[1:v]scale={AV_DISP_W}:{AV_DISP_H}:"
        f"force_original_aspect_ratio=decrease:flags=lanczos,"
        f"pad={AV_DISP_W}:{AV_DISP_H}:(ow-iw)/2:(oh-ih)/2:color=black@0,"
        f"format=yuva420p[avatar_body];"
        f"[0:v]"
        f"drawbox="
        f"x=iw-{AV_DISP_W + PAD + 3}:"
        f"y=ih-{AV_DISP_H + PAD + 3}:"
        f"w={AV_DISP_W + 6}:"
        f"h={AV_DISP_H + 6}:"
        f"color=00BCD4@0.75:t=4,"
        f"{wm_filter}[bg];"
        f"[bg][avatar_body]overlay="
        f"x='W-{AV_DISP_W + PAD}':"
        f"y='H-{AV_DISP_H + PAD}'[v]"
    )

    cmd = [
        _ffmpeg(), "-y",
        "-i", slides_video,
        "-i", avatar_video,
        "-filter_complex", pip_filter,
        "-map", "[v]", "-map", "1:a",
        "-c:v", "libx264", "-b:v", "3500k", "-crf", "17", "-preset", "slow",
        "-c:a", "aac", "-b:a", "192k", "-ar", "44100",
        output_path,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(f"PiP overlay failed: {result.stderr[-600:]}")
    return output_path


# ── Final concat (intro → body → outro) ──────────────────────────────────────

def concat_final_video(
    intro_path:  str,
    body_path:   str,
    outro_path:  str,
    output_path: str,
):
    """Concatenate intro, body, and outro with crossfade transitions."""
    scale_fill = (
        "scale=1920:1080:force_original_aspect_ratio=increase:flags=lanczos,"
        "crop=1920:1080:(iw-1920)/2:0,setsar=1,fps=24"
    )
    scale_pad = (
        "scale=1920:1080:force_original_aspect_ratio=decrease,"
        "pad=1920:1080:(ow-iw)/2:(oh-ih)/2:color=black,setsar=1,fps=24"
    )

    dur_intro = _get_duration(intro_path)
    dur_body  = _get_duration(body_path)
    xf1 = max(dur_intro - 0.5, 0.1)
    xf2 = xf1 + max(dur_body - 0.5, 0.1)

    concat_filter = (
        f"[0:v]{scale_fill}[v0];"
        f"[1:v]{scale_pad}[v1];"
        f"[2:v]{scale_fill}[v2];"
        f"[v0][v1]xfade=transition=fade:duration=0.5:offset={xf1:.3f}[vx1];"
        f"[vx1][v2]xfade=transition=fade:duration=0.5:offset={xf2:.3f}[vout];"
        "[0:a][1:a][2:a]concat=n=3:v=0:a=1[aout]"
    )
    cmd = [
        _ffmpeg(), "-y",
        "-i", intro_path, "-i", body_path, "-i", outro_path,
        "-filter_complex", concat_filter,
        "-map", "[vout]", "-map", "[aout]",
        "-c:v", "libx264", "-b:v", "3500k", "-crf", "17", "-preset", "slow", "-r", "24",
        "-c:a", "aac", "-b:a", "192k", "-ar", "44100",
        output_path,
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(f"Final concat failed: {result.stderr[-800:]}")
    return output_path