"""
pipeline/slide_generator.py
============================
Converts a lecture script into slide JSON, speaker notes, V-style visual
concepts/images, and a PowerPoint deck. This mirrors the edu-genie-upgraded
notebook pipeline while keeping runtime secrets in .env.

FIX (env loading):
  generate_images_from_concepts() previously had a hardcoded OpenRouter API
  key in the notebook source.  That key has been removed; the function now
  requires the caller to pass api_key explicitly (job_worker.py already does
  this via settings.OPENROUTER_API_KEY, which is now correctly loaded from
  .env thanks to the config.py fix).
"""

from __future__ import annotations

import base64
import json
import logging
import os
import re
from copy import deepcopy
from itertools import zip_longest
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import httpx
import torch

logger = logging.getLogger("edugenie.slides")


SLIDE_SYSTEM_PROMPT = """
You are a strict JSON extraction engine. Your task is to convert exactly ONE lecture section into structured slide content suitable for a PowerPoint presentation.

Return ONLY a valid JSON object matching the schema below. No markdown block, no backticks, no markdown formatting, and no conversational text.

### JSON Schema:
{
  "title": "Exact Section Title",
  "content": [
    {
      "point": "Concise bullet point title",
      "explanation": "Supporting details and technical data"
    }
  ]
}

### Extraction Rules:
1. Use ONLY information explicitly present in the provided text.
2. Copy the provided section title exactly into the "title" field.
3. Every paragraph, definition, formula, algorithm, and example must be represented.

### Slide & Content Formatting Rules:
* Points must be concise, 3-8 points per section, maximum 15 words per point.
* Explanations must preserve formulas, definitions, technical metrics, and step-by-step examples.

Return ONLY the raw JSON object.
""".strip()

SPEAKER_NOTE_SYSTEM = """
You are a professor explaining slides to a student. Write like you're talking, not presenting.

Rules:
- 3 to 5 short sentences only
- Each sentence adds something new, never repeats the previous one
- No formal words, no bullet points, just flowing speech

Now write one speaker note for this slide. Output only the note, nothing else.
""".strip()

V_SYSTEM_PROMPT = """
You are an expert educational visual designer and teacher.

TASK:
Convert the given script into multiple visual concepts AND generate a deep teaching explanation for each.

OUTPUT FORMAT:
Return ONLY a valid JSON array matching this schema:
[
  {
    "concept": "",
    "prompt": "",
    "speaker_note": ""
  }
]

VISUAL REQUIREMENTS:
- Generate 3 to 5 visual concepts.
- Each concept represents a different part of the script.
- Each concept must focus on ONE clear idea.

PROMPT RULES:
- Maximum 30 words per prompt.
- Include main objects, environment, and action or relationship.
- Style: clean, educational, easy to understand.

SPEAKER NOTE RULES:
- 3 to 5 sentences.
- Explain the concept in a conversational teaching tone.
- Break the idea into simple, clear parts.
- Do not say "this image shows" or "in this picture".
- Do not introduce new technical facts outside the script.

Return ONLY JSON.
""".strip()

K_SPEAKER_NOTE_SYSTEM = """
You are an expert teacher.
Explain the MCQ answer briefly in a clear teacher tone.
Output only the speaker note.
""".strip()


def _mistral_prompt(system: str, user: str) -> str:
    return f"<s>[INST] {system}\n\n{user} [/INST]"


def _extract_json(text: str) -> Any:
    if not text:
        return None
    text = re.sub(r"```(?:json)?|```", "", text).strip()
    starts = [idx for idx in (text.find("{"), text.find("[")) if idx != -1]
    if not starts:
        return None
    text = text[min(starts):]
    text = re.sub(r",\s*([}\]])", r"\1", text)
    text = re.sub(r'\\([^"\\/bfnrtu])', r"\1", text)
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        for close in ("}", "]"):
            end = text.rfind(close)
            if end != -1:
                try:
                    return json.loads(text[: end + 1])
                except Exception:
                    pass
    return None


def _model_generate(prompt: str, tokenizer, model, max_new_tokens: int, temperature: float = 0.0, device=None) -> str:
    if tokenizer is None or model is None:
        return ""
    inputs = tokenizer(prompt, return_tensors="pt", truncation=True, max_length=2048)
    target_device = device or getattr(model, "device", None)
    if target_device is not None:
        inputs = inputs.to(target_device)
    with torch.no_grad():
        outputs = model.generate(
            **inputs,
            max_new_tokens=max_new_tokens,
            temperature=temperature,
            do_sample=temperature > 0,
            repetition_penalty=1.1,
            eos_token_id=tokenizer.eos_token_id,
            pad_token_id=tokenizer.eos_token_id,
        )
    prompt_len = inputs["input_ids"].shape[1]
    return tokenizer.decode(outputs[0][prompt_len:], skip_special_tokens=True).strip()


def extract_sections(script: str) -> List[dict]:
    """Extract markdown sections and preserve fenced code/equations separately."""
    sections: List[dict] = []
    hierarchy: Dict[int, str] = {}
    current_content: List[str] = []
    code_lines: List[str] = []
    in_code = False

    def clean_title(title: str) -> str:
        title = re.sub(r"^#+\s*", "", title)
        return re.sub(r"\s+", " ", title).strip()

    def save_section():
        nonlocal current_content, code_lines
        content = "\n".join(current_content).strip()
        code = "\n".join(code_lines).strip()
        if not content and not code:
            current_content, code_lines = [], []
            return
        title_parts = [hierarchy[k] for k in sorted(hierarchy)]
        sections.append({
            "title": " - ".join(title_parts) or "Introduction",
            "text": content,
            "content": content,
            "code": code,
        })
        current_content, code_lines = [], []

    for line in script.splitlines():
        if line.strip().startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            code_lines.append(line)
            continue
        match = re.match(r"^(#{1,6})\s+(.*)$", line)
        if match:
            save_section()
            level = len(match.group(1))
            hierarchy[level] = clean_title(match.group(2))
            for key in list(hierarchy):
                if key > level:
                    del hierarchy[key]
            continue
        current_content.append(line)

    save_section()
    if not sections and script.strip():
        sections.append({"title": "Introduction", "text": script.strip(), "content": script.strip(), "code": ""})
    return [s for s in sections if len((s.get("text") or s.get("code") or "").split()) >= 8]


def generate_slide_json(section_title: str, section_text: str, tokenizer, model, device=None, max_new_tokens: int = 1000) -> dict:
    if tokenizer is None or model is None:
        return {"title": section_title, "content": [{"point": section_title, "explanation": section_text[:500]}]}
    user_msg = f"""
Convert the following section into a presentation slide.

Rules:
* Include ALL important content.
* Do not skip information.
* Follow JSON format exactly.

Section_title:
{section_title}

Section_content:
{section_text[:2500]}
""".strip()
    decoded = _model_generate(_mistral_prompt(SLIDE_SYSTEM_PROMPT, user_msg), tokenizer, model, max_new_tokens, device=device)
    slide_data = _extract_json(decoded)
    if not isinstance(slide_data, dict) or "content" not in slide_data:
        logger.warning("Slide JSON parse failed for section '%s'", section_title)
        slide_data = {"title": section_title, "content": [{"point": section_title, "explanation": section_text[:500]}]}
    slide_data["title"] = slide_data.get("title") or section_title
    slide_data["content"] = slide_data.get("content") or []
    return slide_data


def build_slides_list(script: str, tokenizer, model, device=None, max_points: int = 3) -> List[dict]:
    sections = extract_sections(script)
    slides_list: List[dict] = []
    for sec in sections:
        slide_json = generate_slide_json(sec["title"], sec["text"], tokenizer, model, device)
        slide_json["code"] = sec.get("code", "")
        slides_list.append(slide_json)
    return split_large_slides(preprocess_slides_with_code(slides_list), max_points=max_points)


def generate_speaker_note(slide_json: dict, tokenizer, model, device=None) -> str:
    if slide_json.get("_is_k_question"):
        mcq = slide_json.get("_mcq", {})
        return f"Quiz time. {mcq.get('question', '')} Take a moment to think about your answer."
    if slide_json.get("_is_k_answer"):
        mcq = slide_json.get("_mcq", {})
        return (
            f"The correct answer is {mcq.get('correct', '')}. "
            f"{mcq.get('correct_answer', '')}. {mcq.get('explanation', '')}"
        ).strip()

    code = slide_json.get("code") or "\n".join(
        block.get("code", "") for block in slide_json.get("_code_blocks", [])
    )
    code_context = f"\nRelated Code/Equations:\n{code}" if code else ""
    user_msg = f"""
Convert the following presentation slide details into a flowing speaker note.

Section_title:
{slide_json.get("title", "")}

Section_content:
{slide_json.get("content", [])}{code_context}
""".strip()
    return _model_generate(_mistral_prompt(SPEAKER_NOTE_SYSTEM, user_msg), tokenizer, model, 150, device=device)


def generate_visual_concepts(script: str, tokenizer, model, device=None) -> List[dict]:
    decoded = _model_generate(_mistral_prompt(V_SYSTEM_PROMPT, script[:5000]), tokenizer, model, 1500, device=device)
    scenes = _extract_json(decoded)
    if not isinstance(scenes, list):
        logger.warning("V-style visual JSON parse failed")
        return []
    return enhance_prompts(scenes[:5])


def enhance_prompts(visuals_list: List[dict]) -> List[dict]:
    style = "clean educational diagram, 2D illustration, soft colors, high quality"
    enhanced = []
    for item in visuals_list:
        prompt = str(item.get("prompt", "")).strip()
        if not prompt:
            continue
        if len(prompt.split()) > 30:
            prompt = " ".join(prompt.split()[:30])
        enhanced.append({
            "concept": str(item.get("concept", "")).strip(),
            "prompt": f"{prompt}, {style}",
            "speaker_note": str(item.get("speaker_note", "")).strip(),
        })
    return enhanced


def _save_openrouter_image(image_data: str, file_path: Path):
    if image_data.startswith("data:image"):
        img_bytes = base64.b64decode(image_data.split(",", 1)[1])
    else:
        with httpx.Client(timeout=120) as client:
            r = client.get(image_data)
            r.raise_for_status()
            img_bytes = r.content
    file_path.write_bytes(img_bytes)


def generate_images_from_concepts(
    visuals_list: List[dict],
    image_folder: str,
    api_key: str,                                     # FIX: required — no hardcoded key
    model_name: str = "black-forest-labs/flux-1-schnell",
) -> List[str]:
    """
    Generate V-style images through OpenRouter and return saved image paths.

    FIX: api_key is now a required argument.  The notebook had a hardcoded key
    in generate_images_from_concepts(); that key has been removed.  job_worker.py
    already passes settings.OPENROUTER_API_KEY, which is read from .env.
    """
    if not visuals_list:
        return []
    if not api_key:
        logger.warning("OPENROUTER_API_KEY is empty; V-style images will be skipped")
        return []

    out_dir = Path(image_folder)
    out_dir.mkdir(parents=True, exist_ok=True)
    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    saved: List[str] = []

    with httpx.Client(timeout=180) as client:
        for idx, item in enumerate(visuals_list):
            prompt = item.get("prompt", "")
            if not prompt:
                continue
            payload = {
                "model": model_name,
                "messages": [{"role": "user", "content": prompt}],
                "modalities": ["image"],
            }
            try:
                response = client.post(url, headers=headers, json=payload)
                response.raise_for_status()
                result = response.json()
                image_data = result["choices"][0]["message"]["images"][0]["image_url"]["url"]
                file_path = out_dir / f"image_{idx}.png"
                _save_openrouter_image(image_data, file_path)
                item["image_path"] = str(file_path)
                saved.append(str(file_path))
            except Exception as exc:
                logger.error("OpenRouter image generation failed for visual %s: %s", idx, exc)
    return saved


def summarize_one_explanation(explanation: str, tokenizer, model) -> str:
    prompt = f"""Summarize the following explanation into ONE short clear sentence.

Rules:
- Do NOT repeat the input.
- Do NOT include numbering.
- Do NOT include instructions.
- Keep it concise.

Explanation:
{explanation}

Summary:
"""
    summary = _model_generate(prompt, tokenizer, model, 50)
    return summary.split("Summary:")[-1].strip() if summary else explanation[:160]


def summarize_explanations(slides_list: List[dict], tokenizer, model) -> List[dict]:
    for slide in slides_list:
        for item in slide.get("content", []):
            explanation = item.get("explanation", "")
            item["short_explanation"] = summarize_one_explanation(explanation, tokenizer, model) if explanation else ""
    return slides_list


def merge_visual_speaker_notes(slides_list: List[dict], slide_notes: List[str], visuals_list: List[dict]) -> List[str]:
    visual_notes = [item.get("speaker_note", "") for item in visuals_list if item.get("image_path") or item.get("speaker_note")]
    merged: List[str] = []
    for slide, note, visual_note in zip_longest(slides_list, slide_notes, visual_notes):
        if slide is not None and note is not None:
            merged.append(note)
        if visual_note:
            merged.append(visual_note)
    return merged


def wikipedia_search(query: str, limit: int = 5) -> List[dict]:
    """Fetch R-style reference links, matching the notebook's Wikipedia step."""
    _WIKI_HEADERS = {
        "User-Agent": "EduGenieApp/1.0 (educational video generator; contact: edugenie@example.com)",
        "Accept": "application/json",
    }
    try:
        with httpx.Client(timeout=20, headers=_WIKI_HEADERS, follow_redirects=True) as client:
            response = client.get(
                "https://en.wikipedia.org/w/api.php",
                params={
                    "action": "opensearch",
                    "search": query,
                    "limit": limit,
                    "namespace": 0,
                    "format": "json",
                },
            )
            response.raise_for_status()
            data = response.json()
        return [{"title": title, "link": link} for title, link in zip(data[1], data[3])]
    except Exception as exc:
        logger.warning("Wikipedia reference lookup failed: %s", exc)
        return []


def preprocess_slides_with_code(slides_list: List[dict]) -> List[dict]:
    new_slides: List[dict] = []
    for slide in slides_list:
        content = slide.get("content", [])
        text_content = []
        code_blocks = []

        slide_level_code = str(slide.get("code", "")).strip()
        if slide_level_code:
            code_blocks.append({"label": "Source Implementation", "code": slide_level_code})

        for item in content:
            item_code = str(item.get("code", "")).strip()
            if item_code:
                code_blocks.append({"label": item.get("point", ""), "code": item_code})
                text_content.append({**item, "code": ""})
            else:
                text_content.append(item)

        clean_slide = {**slide, "content": text_content}
        clean_slide.pop("code", None)
        new_slides.append(clean_slide)

        if code_blocks:
            new_slides.append({
                "title": f"{slide.get('title', '')} - Code & Equations",
                "content": [],
                "_code_blocks": code_blocks,
                "_is_code_slide": True,
            })
    return new_slides


def split_large_slides(slides_list: List[dict], max_points: int = 3) -> List[dict]:
    new_slides: List[dict] = []
    for slide in slides_list:
        content = slide.get("content", [])
        if slide.get("_is_code_slide") or len(content) <= max_points:
            new_slides.append(slide)
            continue
        for i in range(0, len(content), max_points):
            chunk = content[i:i + max_points]
            part = deepcopy(slide)
            part["content"] = chunk
            part["title"] = f"{slide.get('title', '')} - Part {(i // max_points) + 1}"
            new_slides.append(part)
    return new_slides


def _option_text(options: Any, letter: str, index: int) -> str:
    if isinstance(options, dict):
        return str(options.get(letter, ""))
    if isinstance(options, list) and index < len(options):
        return str(options[index])
    return ""


def inject_k_questions(slides_list: List[dict], mcq_list: List[dict]) -> Tuple[List[dict], List[dict]]:
    if not mcq_list:
        return slides_list, []

    enriched: List[dict] = []
    used: List[dict] = []
    mcq_iter = iter(mcq_list)
    question_freq = max(1, len(slides_list) // max(1, min(len(mcq_list), len(slides_list))))

    for i, slide in enumerate(slides_list):
        enriched.append(slide)
        if (i + 1) % question_freq != 0:
            continue
        try:
            mcq = next(mcq_iter)
        except StopIteration:
            continue
        used.append(mcq)
        opts_text = "\n".join(
            f"{letter}) {_option_text(mcq.get('options', {}), letter, idx)}"
            for idx, letter in enumerate(["A", "B", "C", "D"])
        )
        enriched.append({
            "title": "Quick Check",
            "content": [{"point": mcq.get("question", ""), "explanation": opts_text}],
            "_is_k_question": True,
            "_mcq": mcq,
        })
        enriched.append({
            "title": "Answer & Explanation",
            "content": [{
                "point": f"Correct Answer: {mcq.get('correct', '')}) {mcq.get('correct_answer', '')}",
                "explanation": mcq.get("explanation", ""),
            }],
            "_is_k_answer": True,
            "_mcq": mcq,
        })
    return enriched, used


def _set_slide_notes(slide, note_text: str):
    """
    Write *note_text* into the notes pane of *slide*, exactly as the
    edu-genie-upgraded notebook records speaker notes per slide.
    """
    if not note_text or not note_text.strip():
        return
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = note_text.strip()
    except Exception as exc:
        logger.debug("Could not set notes on slide: %s", exc)


def create_presentation(
    slides_list: List[dict],
    template_path: Optional[str],
    output_path: str,
    learning_style: str = "R",
    visuals_list: Optional[List[dict]] = None,
    results: Optional[list] = None,
    image_path: Optional[str] = None,
    mcq_list: Optional[List[dict]] = None,
    title: str = "Lecture",
    speaker_notes: Optional[List[str]] = None,
) -> str:
    """
    Build the PPTX.

    speaker_notes is a flat list aligned 1-to-1 with the entries in
    slides_list (same order that job_worker.py generates TTS audio).
    Each non-empty string is written into the corresponding slide's notes
    pane so LibreOffice / PowerPoint displays it as a speaker note, exactly
    mirroring the notebook pattern.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("python-pptx not installed - run: pip install python-pptx") from exc

    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
        while len(prs.slides) > 0:
            r_id = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()

    if not slides_list:
        raise ValueError("Slides list is empty")

    LS = learning_style.upper()
    notes_iter = iter(speaker_notes or [])

    # ── Cover slide ────────────────────────────────────────────────────────────
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    if cover.shapes.title:
        cover.shapes.title.text = title
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = "EduGenie Adaptive Learning"

    visual_index = 0
    image_folder = image_path or ""

    for slide_info in slides_list:
        note = next(notes_iter, "")

        if slide_info.get("_is_code_slide"):
            added = _add_code_slide(prs, slide_info)
        else:
            added = _add_content_slide(prs, slide_info, LS)

        _set_slide_notes(added, note)

        if LS == "V" and visuals_list and visual_index < len(visuals_list):
            img = visuals_list[visual_index].get("image_path") or os.path.join(image_folder, f"image_{visual_index}.png")
            if img and os.path.exists(img):
                img_slide = prs.slides.add_slide(prs.slide_layouts[6])
                img_slide.shapes.add_picture(
                    img,
                    (prs.slide_width - Inches(8)) / 2,
                    (prs.slide_height - Inches(5)) / 2,
                    width=Inches(8),
                    height=Inches(5),
                )
                visual_note = next(notes_iter, "") or note
                _set_slide_notes(img_slide, visual_note)
                visual_index += 1

    if LS == "R" and results:
        ref_slide = prs.slides.add_slide(prs.slide_layouts[1])
        if ref_slide.shapes.title:
            ref_slide.shapes.title.text = "Further Reading For You"
        if len(ref_slide.placeholders) > 1:
            tf = ref_slide.placeholders[1].text_frame
            tf.clear()
            for i, ref in enumerate(results[:8]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(ref.get("title", ref))
                link = ref.get("link") if isinstance(ref, dict) else None
                if link:
                    p2 = tf.add_paragraph()
                    run = p2.add_run()
                    run.text = "Open Source"
                    run.hyperlink.address = link
                    p2.level = 1

    prs.save(output_path)
    logger.info("PPTX saved to %s", output_path)
    return output_path


def _add_content_slide(prs, slide_info: dict, learning_style: str):
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = slide_info.get("title", "")
    if len(slide.placeholders) <= 1:
        return slide
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.word_wrap = True
    for j, item in enumerate(slide_info.get("content", [])[:8]):
        point = str(item.get("point", ""))
        explanation = item.get("explanation", "")
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.bold = True
        if explanation:
            text = explanation if learning_style == "R" else item.get("short_explanation", explanation)
            if isinstance(text, dict):
                text = "\n".join(f"{k}: {v}" for k, v in text.items())
            elif isinstance(text, list):
                text = "\n".join(str(x) for x in text)
            p2 = tf.add_paragraph()
            p2.text = str(text)
            p2.level = 1
    return slide


def _add_code_slide(prs, slide_info: dict):
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = slide_info.get("title", "Code & Equations")
    if len(slide.placeholders) <= 1:
        return slide
    ph = slide.placeholders[1]
    left, top, width, height = ph.left, ph.top, ph.width, ph.height
    ph._element.getparent().remove(ph._element)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    text = "\n\n".join(block.get("code", "") for block in slide_info.get("_code_blocks", []))
    p = tf.paragraphs[0]
    p.text = text[:3500]
    p.font.name = "Consolas"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(40, 40, 40)
    return slide